from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # slide background
TEAL        = RGBColor(0x00, 0xB4, 0xD8)   # accent / headings
LIGHT_TEAL  = RGBColor(0x90, 0xE0, 0xEF)   # sub-headings
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xD0, 0xD8, 0xE8)
GOLD        = RGBColor(0xFF, 0xD1, 0x66)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
RED_LIGHT   = RGBColor(0xFF, 0x69, 0x6D)
CARD_BG     = RGBColor(0x14, 0x2A, 0x5A)   # card / box background

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color if line_color else fill_color
    shape.line.width = line_width
    return shape

def accent_bar(slide, t=1.18, color=TEAL):
    r = slide.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.33), Inches(0.055))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def slide_title(slide, title, subtitle=None):
    accent_bar(slide)
    add_text(slide, title, 0.5, 0.22, 12.33, 0.85,
             font_size=28, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.0, 12.33, 0.4,
                 font_size=14, color=LIGHT_TEAL, align=PP_ALIGN.LEFT)

def bullet_box(slide, items, l, t, w, h, title=None,
               font_size=13, title_color=GOLD, bullet="▸ "):
    if title:
        add_text(slide, title, l, t, w, 0.38,
                 font_size=14, bold=True, color=title_color)
        t += 0.38
        h -= 0.38
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = LIGHT_GREY

def badge(slide, text, l, t, w=1.4, h=0.38, bg_color=TEAL, txt_color=NAVY, font_size=11):
    r = rect(slide, l, t, w, h, fill_color=bg_color)
    r.line.fill.background()
    add_text(slide, text, l+0.04, t+0.04, w-0.08, h-0.08,
             font_size=font_size, bold=True, color=txt_color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# big accent line
r = s.shapes.add_shape(1, Inches(0), Inches(3.55), Inches(13.33), Inches(0.07))
r.fill.solid(); r.fill.fore_color.rgb = TEAL; r.line.fill.background()

add_text(s, "NLPRec", 0.7, 0.55, 11.93, 1.1,
         font_size=56, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

add_text(s,
         "AI-Powered Course Recommendation System\nUsing Natural Language Processing",
         0.7, 1.65, 11.93, 1.3,
         font_size=24, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
         "Architecture  ·  Personalization  ·  Evaluation",
         0.7, 2.95, 11.93, 0.55,
         font_size=16, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

add_text(s,
         "Department of Computer Science and Engineering",
         0.7, 3.85, 11.93, 0.45,
         font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(s,
         "First Author · Second Author · Third Author",
         0.7, 4.28, 11.93, 0.4,
         font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(s, "March 2026", 0.7, 6.9, 11.93, 0.4,
         font_size=11, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ABSTRACT / OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Abstract — At a Glance")

rect(s, 0.5, 1.4, 12.33, 5.6, fill_color=CARD_BG)

add_text(s,
         "MOOCs are exploding — 7,000+ courses on Coursera alone. Finding the right course from "
         "a natural-language query like 'I want to learn AI but I am bad at math and a total beginner' "
         "is beyond keyword search.",
         0.75, 1.55, 11.83, 1.1, font_size=14, color=WHITE)

# stat boxes
for (val, lbl, col, lx) in [
    ("0.72",  "Precision@5",  TEAL,       0.75),
    ("0.98",  "Recall@5",     GREEN,       3.45),
    ("0.82",  "F1@5",         GOLD,        6.15),
    ("+71%",  "Avg Δ vs KW",  RED_LIGHT,   8.85),
]:
    rect(s, lx, 2.85, 2.4, 1.5, fill_color=NAVY)
    add_text(s, val, lx, 3.0, 2.4, 0.72,
             font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, lx, 3.72, 2.4, 0.4,
             font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

bullet_box(s,
    ["Seven-stage NLP preprocessing pipeline with negation preservation",
     "Sublinear TF-IDF with bigram features (5K vocabulary)",
     "Nine-step query understanding engine — abbrev expansion, spell correction, difficulty extraction",
     "Adaptive per-user profiling with recency-weighted topic modelling",
     "Live DuckDuckGo search with 24-hour caching & on-the-fly re-ranking",
     "Open-source Streamlit web application"],
    0.75, 4.55, 11.83, 2.3, font_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM & MOTIVATION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Problem Statement & Motivation",
            "Why keyword search fails for course discovery")

# Left card – market
rect(s, 0.5, 1.45, 5.9, 5.6, fill_color=CARD_BG)
add_text(s, "📈 The Scale Problem", 0.7, 1.55, 5.5, 0.45,
         font_size=15, bold=True, color=GOLD)
bullet_box(s,
    ["E-learning market: $250B (2023) → $842B (2030)",
     "Coursera alone: 7,000+ courses, 200+ universities",
     "edX, MIT OCW, Khan Academy growing rapidly",
     "Result: learners face information overload, not scarcity"],
    0.7, 2.05, 5.5, 2.5, font_size=13)

add_text(s, "❌ What Keyword Search Misses", 0.7, 4.55, 5.5, 0.45,
         font_size=13, bold=True, color=RED_LIGHT)
rect(s, 0.7, 5.05, 5.5, 1.7, fill_color=NAVY)
add_text(s,
         '"I want ML — practical, not too mathy, preferably free"',
         0.85, 5.12, 5.2, 0.55,
         font_size=12, italic=True, color=GOLD)
bullet_box(s,
    ["Ignores difficulty constraint",
     "Drops negation: 'not mathy' → ranked by math density",
     "Misses price preference entirely"],
    0.85, 5.65, 5.2, 1.2, font_size=12, bullet="✗ ")

# Right card – gaps
rect(s, 6.65, 1.45, 6.18, 5.6, fill_color=CARD_BG)
add_text(s, "🔍 Three Critical Gaps", 6.85, 1.55, 5.8, 0.45,
         font_size=15, bold=True, color=TEAL)

for (g, desc, cy) in [
    ("Gap 1 — Intent Understanding",
     "Systems discard difficulty signals, negation, and domain abbreviations before retrieval.",
     1.98),
    ("Gap 2 — Cross-Platform Coverage",
     "Existing recommenders trained on single-platform datasets — no multi-source reality.",
     3.38),
    ("Gap 3 — Cold-Start",
     "CF systems fail for new users who make up the bulk of first-time MOOC learners.",
     4.78),
]:
    rect(s, 6.85, cy, 5.8, 1.2, fill_color=NAVY)
    add_text(s, g, 7.0, cy+0.07, 5.5, 0.38,
             font_size=12, bold=True, color=TEAL)
    add_text(s, desc, 7.0, cy+0.45, 5.5, 0.7,
             font_size=11, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — CONTRIBUTIONS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Key Contributions", "Six pillars of NLPRec")

contribs = [
    ("C1", "NLPRec Framework",
     "Open-source modular NLP recommendation pipeline integrating preprocessing, TF-IDF retrieval & engagement amplification",
     TEAL, 0.5, 1.45),
    ("C2", "Query Understanding Engine",
     "Nine-step pipeline: abbreviation expansion, domain-safe spell-correction, difficulty extraction, intent-noise stripping",
     GOLD, 4.05, 1.45),
    ("C3", "Adaptive User Profiling",
     "Per-user JSON profiles with recency-weighted topic preferences learned from implicit behaviour",
     GREEN, 7.6, 1.45),
    ("C4", "Engagement-Augmented Ranking",
     "Log-dampened boost that surfaces popular courses without letting legacy content dominate query-relevant results",
     LIGHT_TEAL, 0.5, 4.05),
    ("C5", "Evaluation Framework",
     "IR-style P@K, R@K, F1@K with fuzzy relevance matching over 10 curated, multi-constraint test queries",
     GOLD, 4.05, 4.05),
    ("C6", "Live Search Integration",
     "Real-time DuckDuckGo search with content filtering, on-the-fly TF-IDF re-ranking & 24-hour disk caching",
     RED_LIGHT, 7.6, 4.05),
]

for (code, name, desc, col, lx, ty) in contribs:
    rect(s, lx, ty, 3.3, 2.35, fill_color=CARD_BG)
    add_text(s, code, lx+0.15, ty+0.1, 0.65, 0.55,
             font_size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, name, lx+0.85, ty+0.12, 2.25, 0.45,
             font_size=12, bold=True, color=WHITE)
    add_text(s, desc, lx+0.15, ty+0.65, 2.95, 1.55,
             font_size=11, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — RELATED WORK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Related Work & Positioning")

# Table header
rect(s, 0.5, 1.42, 12.33, 0.42, fill_color=TEAL)
for (txt, lx, w) in [
    ("System",         0.55, 2.6),
    ("NL Query",       3.2,  1.4),
    ("Cold-Start",     4.65, 1.4),
    ("Multi-Platform", 6.1,  1.85),
    ("Engagement",     7.98, 1.5),
    ("Live Search",    9.53, 1.5),
    ("Open-Source",    11.08,1.55),
]:
    add_text(s, txt, lx, 1.47, w, 0.35,
             font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

rows = [
    ("CF-based [3]",       "✗","✗","✗","✗","✗","✓", CARD_BG),
    ("Ontology CBF [7]",   "~","✓","✗","✗","✗","✗", CARD_BG),
    ("BERT-MOOC [10]",     "✓","✓","✗","✗","✗","✗", CARD_BG),
    ("Wan & Niu [6]",      "✗","✗","✗","✓","✗","✗", CARD_BG),
    ("NLPRec (ours)",      "✓","✓","✓","✓","✓","✓", RGBColor(0x09,0x35,0x26)),
]

for i,(sys,q,cs,mp,eng,ls,os_,rbg) in enumerate(rows):
    ty = 1.87 + i*0.88
    tc = WHITE if rbg != RGBColor(0x09,0x35,0x26) else GREEN
    rect(s, 0.5, ty, 12.33, 0.82, fill_color=rbg)
    add_text(s, sys, 0.55, ty+0.18, 2.6, 0.5,
             font_size=12, bold=(sys=="NLPRec (ours)"), color=tc)
    for (val, lx, w) in [
        (q,   3.2,  1.4),
        (cs,  4.65, 1.4),
        (mp,  6.1,  1.85),
        (eng, 7.98, 1.5),
        (ls,  9.53, 1.5),
        (os_, 11.08,1.55),
    ]:
        col = GREEN if val=="✓" else (GOLD if val=="~" else RED_LIGHT)
        add_text(s, val, lx, ty+0.18, w, 0.5,
                 font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

add_text(s, "NLPRec is the only system ticking all six criteria simultaneously.",
         0.5, 6.65, 12.33, 0.5, font_size=13, italic=True,
         color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "System Architecture", "Eight-phase modular pipeline")

phases = [
    ("1", "Data\nScraping",      "scraper.py",             TEAL),
    ("2", "NLP\nPreprocess",     "text_preprocessing.py",  GOLD),
    ("3", "TF-IDF\nVectorise",   "vectorizer.py",          GREEN),
    ("4", "Similarity\nRanking", "recommender.py",         TEAL),
    ("5", "User\nProfiling",     "user_profile.py",        GOLD),
    ("6", "Behavior\nTracking",  "behavior_tracker.py",    GREEN),
    ("7", "Evaluation",          "evaluation.py",          TEAL),
    ("8", "Streamlit\nFront-End","app.py",                 GOLD),
]

box_w = 1.35
box_h = 1.6
gap   = 0.15
start_l = 0.5

for i,(num,lbl,module,col) in enumerate(phases):
    lx = start_l + i*(box_w+gap)
    rect(s, lx, 1.55, box_w, box_h, fill_color=CARD_BG)
    # phase number
    add_text(s, num, lx, 1.6, box_w, 0.45,
             font_size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, lx, 2.05, box_w, 0.7,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, module, lx, 2.78, box_w, 0.35,
             font_size=8, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)
    # arrow (except last)
    if i < len(phases)-1:
        ax = lx + box_w + 0.02
        add_text(s, "→", ax, 1.95, gap+0.05, 0.5,
                 font_size=16, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Side modules
rect(s, 0.5, 3.45, 5.9, 3.55, fill_color=CARD_BG)
add_text(s, "Supporting Modules", 0.7, 3.52, 5.5, 0.38,
         font_size=14, bold=True, color=TEAL)
bullet_box(s,
    ["query_engine.py  — Nine-step query understanding pipeline",
     "live_search.py   — DuckDuckGo real-time search + re-ranking",
     "query_suggestions.py — 30-topic knowledge-graph chip suggestions"],
    0.7, 3.95, 5.5, 2.7, font_size=12)

rect(s, 6.65, 3.45, 6.18, 3.55, fill_color=CARD_BG)
add_text(s, "Design Principle", 6.85, 3.52, 5.8, 0.38,
         font_size=14, bold=True, color=GOLD)
add_text(s,
         "Each module exposes a clean, typed API — any component (e.g. TF-IDF vectoriser) "
         "can be swapped for a more powerful alternative without touching the rest of the system.",
         6.85, 3.98, 5.8, 1.5, font_size=12, color=LIGHT_GREY)

add_text(s,
         "Data sources: Coursera REST API   ·   edX Discovery API   ·   MIT OCW Sitemap  →  dataset/courses.csv",
         0.5, 7.1, 12.33, 0.35,
         font_size=11, italic=True, color=LIGHT_TEAL, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — NLP PREPROCESSING PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "NLP Preprocessing Pipeline",
            "Identical 7-stage transformation applied to both corpus documents and user queries")

stages = [
    ("1", "Lowercasing",              "Uniform case — 'Python' == 'python'",                     TEAL),
    ("2", "URL Removal",              "Strip https?://\\S+ from scraped descriptions",            GOLD),
    ("3", "Punctuation & Digit Strip","str.translate()  removes noise tokens",                   GREEN),
    ("4", "Whitespace Normalisation", "Collapse multiple spaces, strip leading/trailing",         TEAL),
    ("5", "Tokenisation",             "NLTK word_tokenize: handles contractions & punctuation",   GOLD),
    ("6", "Selective Stopword Removal","Keep: not, no, nor, never, when, where, what, how, me, my, i",GREEN),
    ("7", "Lemmatisation",            "WordNetLemmatizer: algorithms→algorithm, studying→study",  TEAL),
]

for i,(num,name,desc,col) in enumerate(stages):
    ty = 1.42 + i*0.82
    rect(s, 0.5, ty, 12.33, 0.76, fill_color=CARD_BG)
    add_text(s, num, 0.6, ty+0.12, 0.5, 0.52,
             font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, name, 1.2, ty+0.13, 3.5, 0.45,
             font_size=13, bold=True, color=WHITE)
    add_text(s, desc, 4.8, ty+0.18, 7.8, 0.45,
             font_size=12, color=LIGHT_GREY)

# call-out box
rect(s, 0.5, 7.05, 12.33, 0.42, fill_color=RGBColor(0x09,0x35,0x26))
add_text(s,
         "Key insight: same pipeline on both corpus & query guarantees vectors are compared in a uniformly "
         "transformed representation — preventing silent accuracy loss.",
         0.7, 7.08, 11.93, 0.38,
         font_size=11, italic=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — TF-IDF VECTORISATION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "TF-IDF Vectorisation",
            "Sublinear TF · Smoothed IDF · Bigram features")

# Equations
rect(s, 0.5, 1.42, 5.9, 5.6, fill_color=CARD_BG)
add_text(s, "Core Equations", 0.7, 1.52, 5.5, 0.4,
         font_size=15, bold=True, color=TEAL)

eqs = [
    ("Sublinear TF",
     "TF(t,d) = 1+log(count(t,d))  if count>0\n          = 0                  otherwise",
     "Prevents saturation — 50 mentions of 'Python'\nnot ranked 50× higher than 1 mention."),
    ("Smoothed IDF",
     "IDF(t) = log((1+N)/(1+df(t))) + 1",
     "Avoids zero-division; handles unseen query terms."),
    ("TF-IDF Score",
     "TFIDF(t,d) = TF(t,d) × IDF(t)",
     "Final weight for term t in document d."),
    ("Corpus Document",
     "dᵢ = titleᵢ ⊕ descriptionᵢ ⊕ skillsᵢ",
     "Combining three fields gives far richer signal per course."),
]

ty = 1.98
for (lbl,eq,note) in eqs:
    add_text(s, lbl, 0.7, ty, 5.5, 0.32,
             font_size=12, bold=True, color=GOLD)
    add_text(s, eq, 0.7, ty+0.32, 5.5, 0.52,
             font_size=11, color=TEAL)
    add_text(s, note, 0.7, ty+0.84, 5.5, 0.45,
             font_size=10, italic=True, color=LIGHT_GREY)
    ty += 1.35

# Hyperparameters table
rect(s, 6.65, 1.42, 6.18, 5.6, fill_color=CARD_BG)
add_text(s, "Vectoriser Hyperparameters", 6.85, 1.52, 5.8, 0.4,
         font_size=15, bold=True, color=TEAL)

rect(s, 6.85, 1.97, 5.8, 0.38, fill_color=TEAL)
for (txt, lx, w) in [("Parameter",6.9,2.2),("Value",9.15,1.1),("Rationale",10.3,2.1)]:
    add_text(s, txt, lx, 2.02, w, 0.3,
             font_size=11, bold=True, color=NAVY)

params = [
    ("max_features", "5,000",    "Coverage vs. memory balance"),
    ("ngram_range",  "(1, 2)",   "Captures multi-word domain terms"),
    ("min_df",       "1",        "Include all terms (small corpus)"),
    ("sublinear_tf", "True",     "Prevents term saturation"),
    ("smooth_idf",   "True",     "Prevents zero-division"),
    ("norm",         "L2",       "Cosine → fast dot product"),
]

for i,(p,v,r) in enumerate(params):
    ty2 = 2.38 + i*0.76
    bg2 = NAVY if i%2==0 else CARD_BG
    rect(s, 6.85, ty2, 5.8, 0.72, fill_color=bg2)
    add_text(s, p, 6.9, ty2+0.15, 2.2, 0.42,
             font_size=11, bold=True, color=LIGHT_TEAL)
    add_text(s, v, 9.15, ty2+0.15, 1.1, 0.42,
             font_size=11, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, r, 10.3, ty2+0.15, 2.1, 0.42,
             font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — RETRIEVAL & RANKING
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Retrieval & Ranking Model",
            "Cosine similarity + log-dampened engagement boost")

rect(s, 0.5, 1.42, 12.33, 5.6, fill_color=CARD_BG)

# cosine
add_text(s, "Step 1 — Cosine Similarity Retrieval", 0.7, 1.5, 7.5, 0.42,
         font_size=14, bold=True, color=TEAL)
add_text(s,
         "q→ = vectoriser(preprocess(q))\n"
         "sim(q,dᵢ) = (q→ · dᵢ→) / (‖q→‖ · ‖dᵢ→‖)\n"
         "[L2 normalisation → reduces to a fast dot product via sklearn.cosine_similarity()]",
         0.7, 1.95, 11.83, 0.95, font_size=12, color=LIGHT_GREY)

add_text(s, "Step 2 — Log-Dampened Engagement Boost", 0.7, 2.98, 7.5, 0.42,
         font_size=14, bold=True, color=GOLD)
add_text(s,
         "boostᵣₐᵥᵥ(c) = 0.015 × clicks(c) + 0.025 × saves(c)\n"
         "boost(c)     = min( ln(1 + boostᵣₐᵥᵥ(c)) × 0.05 ,  δ=0.12 )",
         0.7, 3.43, 11.83, 0.82, font_size=12, color=LIGHT_GREY)

# three properties
prop_y = 4.32
for (emoji, title, body, lx) in [
    ("📉","Dampening",
     "Log produces diminishing returns — 1,000 clicks ≠ 1,000× benefit", 0.7),
    ("🔒","Bounded",
     "Cap δ=0.12 → popularity cannot shift cosine score by more than 12%", 4.5),
    ("🔖","Save Weight",
     "Saves (0.025) valued 67% higher than clicks (0.015) — stronger intent signal", 8.3),
]:
    rect(s, lx, prop_y, 3.6, 1.45, fill_color=NAVY)
    add_text(s, emoji+" "+title, lx+0.15, prop_y+0.1, 3.3, 0.4,
             font_size=12, bold=True, color=TEAL)
    add_text(s, body, lx+0.15, prop_y+0.52, 3.3, 0.85,
             font_size=11, color=LIGHT_GREY)

add_text(s, "Step 3 — Final Ranking Score", 0.7, 5.85, 7.5, 0.42,
         font_size=14, bold=True, color=GREEN)
add_text(s,
         "sᵢ = sim(q, dᵢ) + boost(dᵢ)        [threshold: sᵢ ≤ 0.05 discarded; ties broken by rating]",
         0.7, 6.3, 11.83, 0.45, font_size=12, color=LIGHT_GREY)

add_text(s, "Keyword Baseline (Eq. 11):", 0.7, 6.82, 4.0, 0.35,
         font_size=12, bold=True, color=RED_LIGHT)
add_text(s, "scoreₖᵥᵥ(q,d) = Σ 𝟙[t ∈ text(d)]   — no weighting, no preprocessing, no difficulty awareness",
         4.75, 6.82, 7.8, 0.35, font_size=12, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — QUERY UNDERSTANDING ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Query Understanding Engine",
            "Nine-step pipeline running before TF-IDF vectorisation")

stages_q = [
    ("1","Punctuation\nNormalisation",  "Non-standard symbols & repeated punctuation → whitespace",      TEAL),
    ("2","Abbreviation\nExpansion",     "100+ regex rules: ml→machine learning, noob→beginner, js→javascript", GOLD),
    ("3","Spell\nCorrection",           "pyspellchecker (edit-dist=1) with 150+ protected tech-vocab terms",GREEN),
    ("4","Difficulty\nExtraction",      "Regex detects beginner / intermediate / advanced → stored as metadata", TEAL),
    ("5","Intent Noise\nStripping",     "Removes: 'I want to learn', 'teach me', 'show me how to'…",     GOLD),
    ("6","Level Word\nRemoval",         "Strips difficulty words from core topic string before retrieval", GREEN),
    ("7","Topic\nExpansion",            "30-topic graph: ml → {scikit-learn, tensorflow, pytorch, …}",   TEAL),
    ("8","Live Query\nGeneration",      "Generates 3-4 enriched DuckDuckGo search strings",              GOLD),
    ("9","Correction\nDisplay",         "Informational message if normalised query differs from raw input",GREEN),
]

bw = 1.3
gap_q = 0.11
sx = 0.38

for i,(num,lbl,desc,col) in enumerate(stages_q):
    lx2 = sx + i*(bw+gap_q)
    rect(s, lx2, 1.4, bw, 1.8, fill_color=CARD_BG)
    add_text(s, num, lx2, 1.47, bw, 0.42,
             font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, lbl, lx2+0.05, 1.88, bw-0.1, 0.62,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages_q)-1:
        add_text(s, "→", lx2+bw, 2.05, gap_q+0.04, 0.4,
                 font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Desc boxes below
for i,(_,_,desc,col) in enumerate(stages_q):
    lx2 = sx + i*(bw+gap_q)
    rect(s, lx2, 3.32, bw, 2.5, fill_color=NAVY)
    add_text(s, desc, lx2+0.07, 3.4, bw-0.14, 2.3,
             font_size=9, color=LIGHT_GREY)

# Spell correction callout
rect(s, 0.38, 5.98, 12.58, 1.4, fill_color=CARD_BG)
add_text(s, "Domain-Protected Spell Correction", 0.6, 6.05, 5.5, 0.4,
         font_size=13, bold=True, color=TEAL)
add_text(s,
         "150+ tech terms whitelisted  (pytorch, sklearn, tensorflow, …)\n"
         "Without protection: pytorch→portrait  |  sklearn→slain  |  numpy→grumpy",
         0.6, 6.48, 12.0, 0.8, font_size=12, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — USER PROFILING & LIVE SEARCH
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "User Profiling & Live Search Integration")

rect(s, 0.5, 1.42, 6.0, 5.6, fill_color=CARD_BG)
add_text(s, "Adaptive User Profiling", 0.7, 1.52, 5.6, 0.4,
         font_size=15, bold=True, color=TEAL)

bullet_box(s,
    ["Per-user JSON profile created on first visit — no rated history needed",
     "Stores: 50 recent searches · 50 bookmarked courses · topic frequency weights · difficulty counts",
     "Recency-Weighted Topics: +1.0 (first encounter) / +0.5 (repeat) — capped at 100 entries",
     "Short-Query Enrichment: 1-word query → +2 profile terms; 3-word → +1 term",
     "Difficulty Auto-Adaptation: argmax over difficulty counts → preferred level applied to future suggestions",
     "Implicit personalisation from observed behaviour — no explicit ratings required"],
    0.7, 2.02, 5.6, 4.85, font_size=12)

rect(s, 6.75, 1.42, 6.08, 5.6, fill_color=CARD_BG)
add_text(s, "Live Search Integration", 6.95, 1.52, 5.7, 0.4,
         font_size=15, bold=True, color=GOLD)

bullet_box(s,
    ["Problem: static corpus grows stale as new courses appear daily"],
    6.95, 2.02, 5.7, 0.55, font_size=12, bullet="❗ ")

add_text(s, "DuckDuckGo Content Filtering:", 6.95, 2.65, 5.7, 0.38,
         font_size=12, bold=True, color=WHITE)
bullet_box(s,
    ["Reject listicles (10 Best… pattern)",
     "Blacklist: reddit, medium, quora, twitter, linkedin",
     "Drop results with no course-describing words in title"],
    6.95, 3.05, 5.7, 1.1, font_size=12, bullet="✗ ")

add_text(s, "On-the-Fly Re-Ranking:", 6.95, 4.22, 5.7, 0.38,
         font_size=12, bold=True, color=WHITE)
add_text(s,
         "simₗᵢᵥₑ(q, rⱼ) = cos(TFIDF(q), TFIDF(titleⱼ ⊕ snippetⱼ))\n"
         "Fresh vectoriser fit on result set — never contaminates corpus vocabulary.",
         6.95, 4.62, 5.7, 0.85, font_size=12, color=LIGHT_GREY)

add_text(s, "24-Hour Disk Caching:", 6.95, 5.55, 5.7, 0.38,
         font_size=12, bold=True, color=WHITE)
add_text(s,
         "cache_key = MD5(normalised_query)\n"
         "Eliminates redundant API round-trips · reduces latency · avoids rate-limit exposure.",
         6.95, 5.95, 5.7, 0.85, font_size=12, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — EVALUATION METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Evaluation Methodology", "IR metrics at K = 5 with fuzzy relevance matching")

rect(s, 0.5, 1.42, 5.9, 5.6, fill_color=CARD_BG)
add_text(s, "Ground-Truth Test Set", 0.7, 1.52, 5.5, 0.38,
         font_size=14, bold=True, color=TEAL)
add_text(s,
         "10 manually chosen queries spanning diverse topics,\n"
         "difficulty levels, and query types (NL constraints,\n"
         "compound technical, abbreviated inputs).\n"
         "Expert review → 2–5 relevant courses per query.",
         0.7, 1.95, 5.5, 1.15, font_size=12, color=LIGHT_GREY)

queries_short = [
    ("Q1","python programming for beginners","4"),
    ("Q2","machine learning no math","5"),
    ("Q3","data science with python & statistics","4"),
    ("Q4","deep learning neural networks advanced","3"),
    ("Q5","web development html css javascript","5"),
    ("Q6","sql database management beginners","3"),
    ("Q7","NLP text analysis","4"),
    ("Q8","cloud devops docker kubernetes","4"),
    ("Q9","linear algebra calculus for ML","3"),
    ("Q10","recommendation systems CF","2"),
]
for i,(q,lbl,rel) in enumerate(queries_short):
    ty2 = 3.18 + i*0.38
    rc = CARD_BG if i%2==0 else NAVY
    rect(s, 0.5, ty2, 5.9, 0.36, fill_color=rc)
    add_text(s, q, 0.65, ty2+0.04, 0.7, 0.3,
             font_size=10, bold=True, color=GOLD)
    add_text(s, lbl, 1.4, ty2+0.04, 4.2, 0.3,
             font_size=10, color=WHITE)
    add_text(s, "Rel:"+rel, 5.05, ty2+0.04, 0.8, 0.3,
             font_size=10, color=GREEN, align=PP_ALIGN.RIGHT)

rect(s, 6.65, 1.42, 6.18, 5.6, fill_color=CARD_BG)
add_text(s, "Metrics & Fuzzy Matching", 6.85, 1.52, 5.8, 0.38,
         font_size=14, bold=True, color=GOLD)

add_text(s, "Fuzzy Relevance Score:", 6.85, 2.0, 5.8, 0.35,
         font_size=12, bold=True, color=WHITE)
add_text(s,
         "m(p,r) = 0.6·Jtoken(p,r) + 0.4·SequenceMatcher(p,r)\n"
         "Threshold θ = 0.55  |  Substring bonus → max(SM, Jtoken, 0.9)",
         6.85, 2.38, 5.8, 0.85, font_size=11, color=LIGHT_GREY)

add_text(s, "IR Metrics at K=5:", 6.85, 3.35, 5.8, 0.35,
         font_size=12, bold=True, color=WHITE)
for (formula, cy) in [
    ("Precision@K  =  |Relevant ∩ Retrieved₁:ₖ| / K",             3.78),
    ("Recall@K     =  |Relevant ∩ Retrieved₁:ₖ| / |Relevant|",    4.22),
    ("F1@K         =  2·P@K·R@K / (P@K + R@K)",                   4.66),
    ("Δₘ           =  (mNLP − mKW) / mKW × 100%",                 5.10),
]:
    rect(s, 6.85, cy, 5.8, 0.4, fill_color=NAVY)
    add_text(s, formula, 6.98, cy+0.06, 5.55, 0.32,
             font_size=11, color=TEAL)

add_text(s, "Future metrics (not yet computed):", 6.85, 5.62, 5.8, 0.35,
         font_size=11, bold=True, color=LIGHT_GREY)
add_text(s,
         "NDCG@K = DCG@K / IDCG@K     ·     MRR = (1/|Q|) Σ 1/rank_q^first",
         6.85, 6.0, 5.8, 0.5, font_size=11, color=LIGHT_TEAL)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Results — NLPRec vs. Keyword Baseline (K = 5)")

# Big headline numbers
for (metric, nlp, kw, delta, col, lx) in [
    ("Precision@5", "0.72", "0.42", "+71.4%", TEAL,      0.5),
    ("Recall@5",    "0.98", "0.57", "+71.9%", GREEN,     3.48),
    ("F1@5",        "0.82", "0.48", "+70.8%", GOLD,      6.46),
]:
    rect(s, lx, 1.4, 2.74, 2.2, fill_color=CARD_BG)
    add_text(s, metric,  lx+0.1, 1.48, 2.54, 0.4,
             font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, nlp,     lx+0.1, 1.88, 2.54, 0.72,
             font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, "NLPRec",lx+0.1, 2.6,  2.54, 0.3,
             font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "KW: "+kw+"  →  Δ "+delta,
             lx+0.1, 2.92, 2.54, 0.42,
             font_size=11, bold=True, color=delta and RED_LIGHT, align=PP_ALIGN.CENTER)

# Per-query table
rect(s, 0.5, 3.78, 12.33, 0.38, fill_color=TEAL)
headers = ["Query","P@5 NLP","P@5 KW","R@5 NLP","R@5 KW","F1@5 NLP","F1@5 KW"]
widths  = [2.1, 1.5, 1.5, 1.55, 1.5, 1.6, 1.55]
lxs = [0.55]
for w in widths[:-1]: lxs.append(lxs[-1]+w)

for j,(h,lx2,w2) in enumerate(zip(headers,lxs,widths)):
    add_text(s, h, lx2, 3.83, w2, 0.3,
             font_size=10, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)

rows_r = [
    ("Q1 — Python basics",  "0.80","0.60","1.00","0.75","0.89","0.67"),
    ("Q2 — ML no math",     "0.80","0.40","0.80","0.40","0.80","0.40"),
    ("Q3 — Data science",   "0.80","0.60","1.00","0.75","0.89","0.67"),
    ("Q4 — Deep learning",  "0.60","0.40","1.00","0.67","0.75","0.50"),
    ("Q5 — Web dev",        "1.00","0.60","1.00","0.60","1.00","0.60"),
    ("Q6 — SQL beginner",   "0.60","0.40","1.00","0.67","0.75","0.50"),
    ("Q7 — NLP analysis",   "0.80","0.40","1.00","0.50","0.89","0.44"),
    ("Q8 — Cloud/DevOps",   "0.80","0.40","1.00","0.50","0.89","0.44"),
    ("Q9 — Math for ML",    "0.60","0.20","1.00","0.33","0.75","0.25"),
    ("Q10 — RecSys/CF",     "0.40","0.20","1.00","0.50","0.57","0.29"),
    ("Mean",                "0.72","0.42","0.98","0.57","0.82","0.48"),
]

for i,row in enumerate(rows_r):
    ty3 = 4.18 + i*0.29
    bg3 = NAVY if i%2==0 else CARD_BG
    if i == len(rows_r)-1: bg3 = RGBColor(0x09,0x35,0x26)
    rect(s, 0.5, ty3, 12.33, 0.28, fill_color=bg3)
    for j,(val,lx2,w2) in enumerate(zip(row,lxs,widths)):
        col3 = WHITE
        if j==0 and i==len(rows_r)-1: col3=GREEN
        elif j in (1,3,5) and i<len(rows_r)-1: col3=TEAL
        elif j in (2,4,6) and i<len(rows_r)-1: col3=RED_LIGHT
        elif j in (1,3,5) and i==len(rows_r)-1: col3=GREEN
        add_text(s, val, lx2, ty3+0.04, w2, 0.22,
                 font_size=9, bold=(i==len(rows_r)-1),
                 color=col3, align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — KEY FINDINGS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Key Findings & Analysis")

findings = [
    (TEAL,  "Multi-Constraint Query Effect",
     ["Q2 (ML no math): keyword 0.40 → NLPRec 0.80 — doubled performance",
      "Difficulty-signal extraction (Stage 4) labels query as 'beginner'",
      "Negation preservation keeps 'no' → TF-IDF matches beginner/no-heavy-math courses"]),
    (GOLD,  "Near-Perfect Recall@5 = 0.98",
     ["9 of 10 queries achieve perfect Recall@5 = 1.00",
      "Only Q2 falls short (0.80) — course uses 'calculus' not 'math' (lexical gap)",
      "Main driver for Future Work F1: dense embedding integration"]),
    (GREEN, "Abbreviation Expansion is Non-Cosmetic",
     ["Without Stage 2: 'ml for beginners' → zero results (token 'ml' not in vocabulary)",
      "Stage 2 expands ml → 'machine learning' before vectorisation",
      "Essential for real learner abbreviated query behaviour"]),
    (RED_LIGHT, "Engagement Boost Ablation",
     ["Disabling engagement boost: no statistically significant change in P@5, R@5, F1@5",
      "Expected — boost refines ranking among equally relevant results",
      "Value is longitudinal: surfaces learner-validated courses over repeated sessions"]),
]

for i,(col,title,points) in enumerate(findings):
    lx3 = 0.5 if i%2==0 else 6.75
    ty3 = 1.42 if i<2 else 4.38
    rect(s, lx3, ty3, 6.08, 2.72, fill_color=CARD_BG)
    add_text(s, title, lx3+0.2, ty3+0.1, 5.68, 0.42,
             font_size=13, bold=True, color=col)
    bullet_box(s, points, lx3+0.2, ty3+0.55, 5.68, 2.0,
               font_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — LIMITATIONS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Limitations", "Honest assessment of current boundaries")

lims = [
    (RED_LIGHT,"Lexical Gap",
     "TF-IDF cannot equate surface-dissimilar synonyms (calculus ≠ math). Dense embeddings (F1) needed."),
    (GOLD,"Small Ground Truth",
     "10 hand-crafted queries cannot represent the full distribution of real learner intent. Larger crowd-sourced set required."),
    (TEAL,"New-Course Cold Start",
     "Courses newly added to corpus carry boost=0 and may underperform against legacy well-clicked alternatives."),
    (GREEN,"Corpus Freshness",
     "Live search fills query-time gaps, but the TF-IDF model itself is not retrained — cached live results may lag days behind new courses."),
    (RED_LIGHT,"English-Only",
     "Preprocessing, stopword list, spell-corrector, and vocabulary are English-exclusive. Multilingual support is future work (F5)."),
]

for i,(col,title,body) in enumerate(lims):
    ty4 = 1.42 + i * 1.1
    rect(s, 0.5, ty4, 12.33, 1.0, fill_color=CARD_BG)
    add_text(s, "⚠ "+title, 0.7, ty4+0.1, 3.5, 0.42,
             font_size=13, bold=True, color=col)
    add_text(s, body, 4.3, ty4+0.13, 8.3, 0.72,
             font_size=12, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — CONCLUSION & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_title(s, "Conclusion & Future Work")

rect(s, 0.5, 1.42, 5.9, 5.6, fill_color=CARD_BG)
add_text(s, "Conclusion", 0.7, 1.52, 5.5, 0.38,
         font_size=15, bold=True, color=TEAL)
bullet_box(s,
    ["NLPRec: end-to-end intelligent course recommendation understanding query intent, not keyword counts",
     "7-stage NLP pipeline · sublinear TF-IDF bigrams · 9-step query engine",
     "Log-dampened engagement boost · recency-weighted user profiling · live search",
     "K=5 evaluation: P@5 +71.4% · R@5 +71.9% · F1@5 +70.8% over keyword baseline",
     "Recall@5 = 0.98 → learners find the right course on first attempt",
     "Open-source Streamlit app — reproducible reference implementation for EdTech"],
    0.7, 1.98, 5.5, 4.9, font_size=12)

rect(s, 6.65, 1.42, 6.18, 5.6, fill_color=CARD_BG)
add_text(s, "Future Work Roadmap", 6.85, 1.52, 5.8, 0.38,
         font_size=15, bold=True, color=GOLD)

fw = [
    ("F1","Dense Embeddings",       "Sentence-BERT / E5 to bridge lexical gap",                      TEAL),
    ("F2","Hybrid Neural-TF-IDF",   "Late fusion of sparse + dense for precision + recall",          GOLD),
    ("F3","Learning Path Generation","Sequential path planning with prerequisite modelling",           GREEN),
    ("F4","Temporal Signals",       "Recency-weighted ranking to reduce stale content",              TEAL),
    ("F5","Multilingual Support",   "mBERT / XLM-R for non-English queries",                        GOLD),
    ("F6","Large-Scale A/B Study",  "Real learners · task completion · retention · satisfaction",    GREEN),
    ("F7","CF Integration",         "Hybrid recommender for users with sufficient history",           TEAL),
    ("F8","NDCG & MRR",             "Rank-sensitive metrics for position-aware evaluation",           GOLD),
]
for i,(code,name,desc,col) in enumerate(fw):
    ty5 = 2.0 + i*0.62
    rect(s, 6.85, ty5, 5.8, 0.57, fill_color=NAVY)
    add_text(s, code, 6.92, ty5+0.1, 0.5, 0.38,
             font_size=11, bold=True, color=col)
    add_text(s, name, 7.48, ty5+0.05, 2.3, 0.28,
             font_size=11, bold=True, color=WHITE)
    add_text(s, desc, 7.48, ty5+0.32, 4.0, 0.26,
             font_size=10, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — THANK YOU / QA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

r2 = s.shapes.add_shape(1, Inches(0), Inches(3.45), Inches(13.33), Inches(0.08))
r2.fill.solid(); r2.fill.fore_color.rgb = TEAL; r2.line.fill.background()

add_text(s, "Thank You", 0.7, 0.6, 11.93, 1.1,
         font_size=58, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, "Questions & Discussion", 0.7, 1.75, 11.93, 0.7,
         font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "NLPRec — Open-Source Intelligent Course Recommendation System",
         0.7, 2.6, 11.93, 0.5, font_size=14, italic=True,
         color=LIGHT_TEAL, align=PP_ALIGN.CENTER)

add_text(s,
         "P@5 = 0.72   ·   R@5 = 0.98   ·   F1@5 = 0.82   ·   Avg Δ > 71% over keyword baseline",
         0.7, 3.72, 11.93, 0.5,
         font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_text(s,
         "Built with: Python · scikit-learn · NLTK · Streamlit · DuckDuckGo Search · pyspellchecker · matplotlib",
         0.7, 4.45, 11.93, 0.45, font_size=12,
         color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(s,
         "Data sources: Coursera API  ·  edX Discovery API  ·  MIT OpenCourseWare Sitemap",
         0.7, 5.0, 11.93, 0.4, font_size=12,
         color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text(s, "Department of Computer Science and Engineering  |  March 2026",
         0.7, 6.8, 11.93, 0.4, font_size=11,
         color=LIGHT_TEAL, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/prathmeshd/Desktop/NLPRec/research_paper/NLPRec_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
print(f"Slides: {len(prs.slides)}")
